"""
Implement your retrieval system here.

The test_files/ folder contains mixed-format documents:
  - PDFs, PowerPoint, Word, emails, JSON glossary

Your retrieve() function should handle both:
  - Broad queries: "What was the research methodology?"
  - Narrow queries: "What does 'frm_brand_awareness' measure?"
"""

import os 
import json 
import email 
from pathlib import Path 
from typing import List, Dict
import numpy as np 
from sentence_transformers import SentenceTransformer #lokales Embedding Model

MODEL = None 
CHUNKS = None
EMBEDDINGS = None #global speichern, damit sie nicht bei jedem Abruf neu berechnet werden müssen

def load_text_from_file(file_path: Path) -> str:
    
    suffix = file_path.suffix.lower()

    if suffix == '.pdf':
        import pdfplumber
        text = ""

        with pdfplumber.open(file_path) as pdf:
            for page in pdf.pages:
                page_text = page.extract_text()

                if page_text:
                    text += page_text + "\n"
        return text
    
    elif suffix == ".pptx":
        from pptx import Presentation 
        text = ""

        prs = Presentation(file_path)
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text+=shape.text +'\n'
        return text
    
    elif suffix == ".docx":
        from docx import Document
        doc = Document(file_path)
        text = ""

        for para in doc.paragraphs:
            text += para.text +"\n"

        return text
    
    elif suffix == ".eml":
        with open(file_path, 'r', encoding = 'utf-8', errors='ignore') as f:
            msg = email.message_from_file(f)

            text=''

            if msg.is_multipart():
                for part in msg.walk():
                    if part.get_content_type() == "text/plain":
                        text += part.get_payload(decode=True).decode('utf-8', errors='ignore')

            else:
                text += msg.get_payload(decode=True).decode('utf-8', errors='ignore')

        return text
    
    elif suffix == ".json":
        with open(file_path, 'r', encoding='utf-8') as f:
            data=json.load(f)

          # jeden Eintrag als eigene Zeile

        if isinstance(data, dict):
            text_parts = []
            for key, value in data.items():
                entry = f"{key}: {json.dumps(value)}"
                text_parts.append(entry)

            return "\n\n".join(text_parts)

        return json.dumps(data, indent=2)

def split_into_chunks(text:str, source:str, chunk_size: int =500, overlap : int = 50) -> List[Dict]:
    """"
    Split text into chunks.
    """
    chunks = []

    if not text.strip():
        return chunks
    
    words = text.split()

    step = chunk_size - overlap

    for i in range(0, len(words), step):
        chunk_words = words[i:i + chunk_size]
        chunk_text = " ".join(chunk_words)

        if len(chunk_text) < 10:
            continue
        chunks.append({
            "text":chunk_text, 
            "source": source
        })
    return chunks
            
def get_embedding(model : SentenceTransformer, text:str) -> np.ndarray:
    """
    Get embedding for a given text using the provided model.
    """
    return model.encode(text, normalize_embeddings=True)

def cosine_similarity(a: np.ndarray, b:np.ndarray)-> float:
    """
    Compute cosine similarity between two vectors.
    """
    return float(np.dot(a, b))

def hybrid_score(query: str, chunk_text: str, 
                 embedding_score: float) -> float:
    
    query_words = set(query.lower().split())
    chunk_words = set(chunk_text.lower().split())
    
    stopwords = {"what", "is", "the", "a", "an", "of", 
                 "does", "mean", "define", "as", "used", "in"}
    query_words -= stopwords
    
    if not query_words:
        return embedding_score
    
    overlap = len(query_words & chunk_words)
    keyword_score = overlap / len(query_words)
    
    return 0.7 * embedding_score + 0.3 * keyword_score

def retrieve(query: str) -> list[dict]:
    """
    Return top-5 most relevant passages for the query.

    Each result must have:
      - "text": str       — the passage content
      - "source": str     — source filename or "glossary"
      - "score": float    — relevance score (higher = better)
    """
    # TODO: implement your retrieval system

    global MODEL, CHUNKS, EMBEDDINGS

    # Model laden 
    if MODEL is None:
        MODEL = SentenceTransformer("all-MiniLM-L6-v2")

    if CHUNKS is None:
        CHUNKS = []

        test_dir = Path("test_files")

        for file_path in test_dir.iterdir():
            if file_path.suffix.lower() not in [".pdf", ".pptx", ".docx", ".eml", ".json"]:
                continue
            
            text = load_text_from_file(file_path)

            if not text.strip():
                continue
            
            if file_path.suffix.lower() == ".json":
                file_chunks = split_into_chunks(
                    text,
                    file_path.name,
                    chunk_size=30,
                    overlap=5
                                )
            else:
                file_chunks = split_into_chunks(
                    text,
                    file_path.name,
                    chunk_size=200,
                    overlap=30)
            CHUNKS.extend(file_chunks)

    if EMBEDDINGS is None:
        texts = [chunk["text"] for chunk in CHUNKS]
        EMBEDDINGS = MODEL.encode(
            texts, 
            normalize_embeddings=True,
            show_progress_bar=True, 
        )

    query_embedding = get_embedding(MODEL, query)

    scores = []

    for i, chunk in enumerate(EMBEDDINGS):
        embedding_score = cosine_similarity(query_embedding, chunk)
        final_score = hybrid_score(query, CHUNKS[i]["text"], embedding_score)
        scores.append((final_score, i))

    scores.sort(reverse=True)
    top5 = scores[:5]

    results=[]

    for score, idx in top5:
        results.append({
            "text" : CHUNKS[idx]["text"],
            "source": CHUNKS[idx]["source"],
            "score": score
        })
    #print(results)

    return results

    



